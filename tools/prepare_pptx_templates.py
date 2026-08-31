"""
Prepara os templates PPTX de cartao de bloqueio (one-time tool).

Para cada PPTX original em 'MODELOS CARTAO BLOQUEIO/' cria uma copia preparada em
'templates/cards/pptx/' com:
- slides extras removidos (quando aplicavel)
- textos de amostra substituidos por tokens {{CAMPO}}
- shapes renomeados CARD{slot}_{CAMPO} / CARD{slot}_FOTO
- textos residuais limpos (ALTEC PEQUENO) e shapes [FOTO] removidos (CSN)
- card.json companheiro com zonas de recorte (fracao da pagina)

Originais NAO sao modificados.

Uso:  python tools/prepare_pptx_templates.py [--previews]
"""
import copy
import io
import json
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Cm, Pt

ROOT = Path(__file__).resolve().parents[1]
SRC_DIR = next(d for d in ROOT.iterdir() if d.is_dir() and d.name.upper().startswith("MODELOS"))
OUT_DIR = ROOT / "templates" / "cards" / "pptx"
EMU_PER_CM = 360000.0


# ── Utilidades ──────────────────────────────────────────────

def walk(shapes, ancestors=()):
    """Itera (shape, ancestors) recursivamente por grupos."""
    for shp in shapes:
        yield shp, ancestors
        if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from walk(shp.shapes, ancestors + (shp,))


def _group_map(g):
    x = g._element.grpSpPr.xfrm
    off = (x.off.x, x.off.y)
    ext = (x.ext.cx, x.ext.cy)
    cho = (x.chOff.x, x.chOff.y)
    che = (x.chExt.cx, x.chExt.cy)
    return off, ext, cho, che


def to_abs(px, py, ancestors):
    """Converte coordenadas do child-space para absolutas do slide (EMU)."""
    for g in reversed(ancestors):
        off, ext, cho, che = _group_map(g)
        sx = ext[0] / che[0] if che[0] else 1.0
        sy = ext[1] / che[1] if che[1] else 1.0
        px = off[0] + (px - cho[0]) * sx
        py = off[1] + (py - cho[1]) * sy
    return px, py


def abs_cm(shape, ancestors):
    ax, ay = to_abs(shape.left or 0, shape.top or 0, ancestors)
    return ax / EMU_PER_CM, ay / EMU_PER_CM


def drop_slide(prs, slide):
    id_lst = prs.slides._sldIdLst
    for sld_id in list(id_lst):
        if prs.part.related_part(sld_id.rId) == slide.part:
            prs.part.drop_rel(sld_id.rId)
            id_lst.remove(sld_id)


def drop_shape(shape):
    el = shape._element
    el.getparent().remove(el)


def set_text_keep_style(text_frame, lines):
    """Reescreve o texto (1 paragrafo por linha) preservando estilo do 1o run de cada para."""
    from pptx.oxml.ns import qn

    tx_body = text_frame._txBody
    paragraphs = text_frame.paragraphs
    while len(paragraphs) < len(lines):
        newp = copy.deepcopy(paragraphs[0]._p)
        tx_body.append(newp)
        paragraphs = text_frame.paragraphs
    while len(paragraphs) > len(lines):
        tx_body.remove(paragraphs[-1]._p)
        paragraphs = text_frame.paragraphs
    for para, line in zip(paragraphs, lines):
        runs = para.runs
        if runs:
            runs[0].text = line
            for r in runs[1:]:
                r.text = ""
        else:
            para.add_run().text = line


def clear_text(text_frame):
    for para in text_frame.paragraphs:
        for r in para.runs:
            r.text = ""


def slot_band(value_cm, bands):
    for k, (lo, hi) in enumerate(bands, start=1):
        if lo <= value_cm < hi:
            return k
    return None


def cm_fractions(bands_x_cm, bands_y_cm, page_w_cm, page_h_cm):
    """Zonas row-major a partir de bandas (cm)."""
    zones = []
    for lo_y, hi_y in bands_y_cm:
        for lo_x, hi_x in bands_x_cm:
            zones.append([
                round(lo_x / page_w_cm, 4), round(lo_y / page_h_cm, 4),
                round(hi_x / page_w_cm, 4), round(hi_y / page_h_cm, 4),
            ])
    return zones


# ── ARCELORMITTAL ───────────────────────────────────────────

def prep_arcelor(prs):
    for s in list(prs.slides)[1:]:
        drop_slide(prs, s)
    slide = prs.slides[0]
    bands_t = [(0, 6.12), (6.12, 12.52), (12.52, 18.95), (18.95, 26.0)]
    report = {"nome": [], "matricula": [], "foto": []}
    for shape, anc in walk(slide.shapes):
        txt = shape.text_frame.text if getattr(shape, "has_text_frame", False) else ""
        l_cm, t_cm = abs_cm(shape, anc)
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            w = (shape.width or 0) / EMU_PER_CM
            h = (shape.height or 0) / EMU_PER_CM
            if 1.8 <= w <= 3.0 and h >= 2.5:
                k = slot_band(t_cm, bands_t)
                if k:
                    shape.name = f"CARD{k}_FOTO"
                    report["foto"].append((k, shape.name))
        elif re.match(r"^\s*matricula\s*:", txt.strip(), re.I):
            k = slot_band(t_cm, bands_t)
            if k:
                set_text_keep_style(shape.text_frame, ["MATRICULA: {{MATRICULA}}"])
                shape.name = f"CARD{k}_MATRICULA"
                report["matricula"].append(k)
        elif re.match(r"^\s*nome\s*:", txt.strip(), re.I):
            k = slot_band(t_cm, bands_t)
            if k:
                label = "NOME" if txt.strip().upper().startswith("NOME:") else "Nome"
                set_text_keep_style(shape.text_frame, [f"{label}: {{{{NOME}}}}"])
                shape.name = f"CARD{k}_NOME"
                report["nome"].append(k)
    zones = cm_fractions(
        [(0.0, 19.05)], bands_t, 19.05, 25.4
    )
    return {"zones": zones, "cards_per_slide": 4, "report": report}


# ── ALTEC PEQUENO ───────────────────────────────────────────

ALTEC_P_FIELDS = {
    # (l_cm, t_cm): (slot, campo)  — posicoes dos shapes de dados
    (2.12, 6.98): (1, "NOME"), (8.18, 6.99): (2, "NOME"), (12.9, 7.06): (3, "NOME"), (18.12, 7.13): (4, "NOME"),
    (1.99, 15.61): (5, "NOME"), (8.43, 15.62): (6, "NOME"), (13.07, 15.43): (7, "NOME"),
    (2.13, 7.47): (1, "FUNCAO"), (7.53, 7.48): (2, "FUNCAO"), (13.0, 7.47): (3, "FUNCAO"), (17.93, 7.66): (4, "FUNCAO"),
    (2.13, 16.22): (5, "FUNCAO"), (7.42, 16.18): (6, "FUNCAO"), (13.02, 16.1): (7, "FUNCAO"),
    (2.68, 8.03): (1, "TELEFONE"), (7.85, 8.04): (2, "TELEFONE"), (13.07, 8.05): (3, "TELEFONE"), (18.54, 8.09): (4, "TELEFONE"),
    (2.76, 16.83): (5, "TELEFONE"), (7.94, 16.74): (6, "TELEFONE"), (13.17, 16.71): (7, "TELEFONE"),
}


def _find_field(l_cm, t_cm):
    for (fl, ft), val in ALTEC_P_FIELDS.items():
        if abs(l_cm - fl) < 0.45 and abs(t_cm - ft) < 0.45:
            return val
    return None


# posicao da foto 3x4 dentro de cada cartao (card_left cm): (slot, left_cm, top_cm)
ALTEC_P_FOTO_POS = {
    1: (2.18 + 1.55, 2.45),   # A (linha 1)
    2: (7.45 + 1.55, 2.45),   # B
    3: (12.71 + 1.55, 2.45),  # C
    4: (18.0 + 1.55, 2.45),   # D
    5: (2.1 + 1.55, 10.92),   # E (linha 2)
    6: (7.45 + 1.55, 10.92),  # F
    7: (12.79 + 1.55, 10.92), # G
}
ALTEC_P_FOTO_W_CM = 1.9
ALTEC_P_FOTO_H_CM = 2.55

# fonte reduzida: caixas de 2,95-5,34cm nao comportam 12pt com nomes reais
ALTEC_P_FONT_PT = {"NOME": 8.0, "FUNCAO": 7.5, "TELEFONE": 8.0}


def _foto_placeholder_jpeg() -> bytes:
    from PIL import Image

    img = Image.new("RGB", (150, 200), (230, 230, 230))
    out = io.BytesIO()
    img.save(out, format="JPEG")
    return out.getvalue()


def prep_altec_pequeno(prs):
    slide = prs.slides[0]
    report = []
    for shape, anc in walk(slide.shapes):
        if not getattr(shape, "has_text_frame", False):
            continue
        txt = shape.text_frame.text
        if not txt.strip():
            continue
        l_cm, t_cm = abs_cm(shape, anc)
        w = (shape.width or 0) / EMU_PER_CM
        h = (shape.height or 0) / EMU_PER_CM
        # retangulo de fundo do cartao (5x8cm): limpa texto residual
        if abs(w - 5.0) < 0.35 and abs(h - 8.0) < 0.35:
            clear_text(shape.text_frame)
            continue
        found = _find_field(l_cm, t_cm)
        if not found:
            continue
        k, campo = found
        set_text_keep_style(shape.text_frame, [f"{{{{{campo}}}}}"])
        shape.name = f"CARD{k}_{campo}"
        for para in shape.text_frame.paragraphs:
            for r in para.runs:
                r.font.size = Pt(ALTEC_P_FONT_PT[campo])
        report.append((k, campo))

    # foto 3x4: o template original nao tem foto — insere placeholder centralizado
    # na area livre entre os logos do topo e a faixa "MANUTENCAO / NAO LIGUE"
    foto_bytes = _foto_placeholder_jpeg()
    for k, (l_cm, t_cm) in ALTEC_P_FOTO_POS.items():
        pic = slide.shapes.add_picture(
            io.BytesIO(foto_bytes),
            Cm(l_cm), Cm(t_cm), Cm(ALTEC_P_FOTO_W_CM), Cm(ALTEC_P_FOTO_H_CM),
        )
        pic.name = f"CARD{k}_FOTO"
        report.append((k, "FOTO"))
    bands_x = [(0, 7.3), (7.3, 12.5), (12.5, 17.82), (17.82, 25.4)]
    bands_y = [(0, 9.37), (9.37, 19.05)]
    # 7 slots usados (o 8o cartao e reserva em branco no original)
    zones = cm_fractions(bands_x, bands_y, 25.4, 19.05)[:7]
    return {"zones": zones, "cards_per_slide": 7, "report": report}


# ── CSN ─────────────────────────────────────────────────────

CSN_FOTOS = {(2.35, 3.73): 1, (14.68, 3.76): 2, (2.2, 13.07): 3, (14.48, 13.0): 4}
CSN_NOMES = {(0.58, 7.07): 1, (12.96, 7.06): 2, (0.66, 16.4): 3, (13.08, 16.51): 4}


def _by_pos(l_cm, t_cm, table, tol=0.6):
    for (fl, ft), val in table.items():
        if abs(l_cm - fl) < tol and abs(t_cm - ft) < tol:
            return val
    return None


def prep_csn(prs):
    slide = prs.slides[0]
    report = []
    for shape, anc in list(walk(slide.shapes)):
        txt = shape.text_frame.text if getattr(shape, "has_text_frame", False) else ""
        l_cm, t_cm = abs_cm(shape, anc)
        if txt.strip() == "[FOTO]":
            drop_shape(shape)
            continue
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            k = _by_pos(l_cm, t_cm, CSN_FOTOS)
            if k:
                shape.name = f"CARD{k}_FOTO"
                report.append((k, "FOTO"))
        elif txt.strip().upper() == "LIDERADO":
            k = 1 if (t_cm < 9.6 and l_cm < 12.8) else 2 if t_cm < 9.6 else 3 if l_cm < 12.8 else 4
            set_text_keep_style(shape.text_frame, ["{{PAPEL}}"])
            shape.name = f"CARD{k}_PAPEL"
            report.append((k, "PAPEL"))
        elif re.match(r"^\s*nome\s*", txt.strip(), re.I):
            k = _by_pos(l_cm, t_cm, CSN_NOMES)
            if k:
                set_text_keep_style(shape.text_frame, ["Nome: {{NOME}}", "Dpto: {{SETOR}}", "Empresa: {{EMPRESA}}"])
                shape.name = f"CARD{k}_DADOS"
                report.append((k, "DADOS"))
    bands_x = [(0, 12.8), (12.8, 25.4)]
    bands_y = [(0, 9.6), (9.6, 19.05)]
    zones = cm_fractions(bands_x, bands_y, 25.4, 19.05)
    return {"zones": zones, "cards_per_slide": 4, "report": report}


# ── LOTOTO ──────────────────────────────────────────────────

def prep_lototo(prs):
    report = []
    for slide in prs.slides:
        for shape, anc in walk(slide.shapes):
            if shape.name == "CaixaDeTexto 7":
                set_text_keep_style(shape.text_frame, ["{{NOME}}"])
                shape.name = "CARD1_NOME"
                report.append("NOME")
            elif shape.name == "CaixaDeTexto 8":
                set_text_keep_style(shape.text_frame, ["{{FUNCAO}}"])
                shape.name = "CARD1_FUNCAO"
                report.append("FUNCAO")
            elif shape.name == "CaixaDeTexto 9":
                set_text_keep_style(shape.text_frame, ["{{MATRICULA}}"])
                shape.name = "CARD1_MATRICULA"
                report.append("MATRICULA")
            elif shape.name == "CaixaDeTexto 15":
                set_text_keep_style(shape.text_frame, ["{{SETOR}}"])
                shape.name = "CARD1_SETOR"
                report.append("SETOR")
            elif shape.name in ("CaixaDeTexto 11", "CaixaDeTexto 16"):
                set_text_keep_style(shape.text_frame, ["{{PAPEL}}"])
                shape.name = "CARD1_PAPEL" if shape.name == "CaixaDeTexto 11" else "CARD1_PAPEL2"
                report.append("PAPEL")
            elif shape.name == "Imagem 19":
                shape.name = "CARD1_FOTO"
                report.append("FOTO")
    return {"zones": [[0.0, 0.0, 1.0, 1.0]], "cards_per_slide": 1, "report": report}


# ── Main ────────────────────────────────────────────────────

TEMPLATES = [
    {
        "src": "MODELO ARCELORMITTAL.pptx",
        "pptx": "ARCELORMITTAL.pptx",
        "card_code": "ARCELORMITTAL",
        "cliente_nome": "ArcelorMittal Tubarao",
        "descricao": "Cartao de bloqueio ArcelorMittal (4 por folha) — PPTX",
        "empresa_default": "ALTEC",
        "prep": prep_arcelor,
    },
    {
        "src": "MODELO CARTÃO ALTEC PEQUENO.pptx",
        "pptx": "ALTEC_PEQUENO.pptx",
        "card_code": "ALTEC-PEQUENO",
        "cliente_nome": "ALTEC (cartao pequeno)",
        "descricao": "Cartao de bloqueio ALTEC pequeno (7 por folha) — PPTX",
        "empresa_default": "ALTEC",
        "prep": prep_altec_pequeno,
    },
    {
        "src": "MODELO CARTÃO DE BLOQUEIO CSN.pptx",
        "pptx": "CSN.pptx",
        "card_code": "CSN",
        "cliente_nome": "CSN",
        "descricao": "Cartao de bloqueio CSN (4 por folha, Lider/Liderado) — PPTX",
        "empresa_default": "ALTEC",
        "prep": prep_csn,
    },
    {
        "src": "MODELO CARTÃO LOTOTO.pptx",
        "pptx": "LOTOTO.pptx",
        "card_code": "LOTOTO",
        "cliente_nome": "LOTO/TO",
        "descricao": "Cartao LOTOTO (1 por folha, Lider/Liderado) — PPTX",
        "empresa_default": "ALTEC",
        "prep": prep_lototo,
    },
]


def main():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    ok = True
    for cfg in TEMPLATES:
        # busca por nome tolerante a acentos
        src = None
        for f in SRC_DIR.glob("*.pptx"):
            if f.name.upper() == cfg["src"].upper() or _norm(f.name) == _norm(cfg["src"]):
                src = f
                break
        if not src:
            print(f"[ERRO] origem nao encontrada: {cfg['src']}")
            ok = False
            continue

        prs = Presentation(str(src))
        result = cfg["prep"](prs)
        out_pptx = OUT_DIR / cfg["pptx"]
        prs.save(str(out_pptx))

        prs2 = Presentation(str(out_pptx))
        card = {
            "card_code": cfg["card_code"],
            "cliente_nome": cfg["cliente_nome"],
            "descricao": cfg["descricao"],
            "template_type": "pptx",
            "pptx_file": cfg["pptx"],
            "cards_per_slide": result["cards_per_slide"],
            "slides": len(prs2.slides),
            "empresa_default": cfg["empresa_default"],
            "card_zones_fraction": result["zones"],
        }
        card_json = OUT_DIR / f"{cfg['card_code']}.card.json"
        card_json.write_text(json.dumps(card, ensure_ascii=False, indent=2), encoding="utf-8")

        n_shapes = sum(1 for _ in walk_all(prs2))
        print(f"[OK] {cfg['card_code']}: slides={len(prs2.slides)} slots={result['cards_per_slide']} "
              f"zonas={len(result['zones'])} shapes={n_shapes}")
        print(f"     relatorio: {result['report']}")
    return 0 if ok else 1


def _norm(name):
    import unicodedata
    return unicodedata.normalize("NFKD", name).encode("ascii", "ignore").decode().upper()


def walk_all(prs):
    for slide in prs.slides:
        yield from walk(slide.shapes)


if __name__ == "__main__":
    sys.exit(main())
