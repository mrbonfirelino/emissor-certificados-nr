"""
Servico de geracao de cartoes de bloqueio via templates PPTX.

Fluxo:
1. Templates ficam em templates/cards/pptx/ (arquivo .pptx + .card.json companheiro)
2. O .card.json declara: card_code, template_type="pptx", pptx_file, cards_per_slide,
   empresa_default e card_zones_fraction (bbox de cada cartao, em fracao da pagina,
   usado no modo "1 cartao por pagina").
3. Shapes do template sao nomeados CARD{slot}_{CAMPO} (ex: CARD1_NOME, CARD2_FOTO).
   Textos contem tokens {{NOME}} {{FUNCAO}} {{TELEFONE}} {{CPF}} {{MATRICULA}}
   {{SETOR}} {{EMPRESA}} {{PAPEL}} — so os presentes no template sao substituidos.
4. Foto: shape Picture nomeado CARD{k}_FOTO — a imagem e trocada mantendo
   posicao/tamanho (blob swap), com center-crop 3x4 para nao distorcer.
5. Conversao PDF: Microsoft PowerPoint via COM (comtypes), UMA sessao por lote.
6. Merge de PDFs via PyMuPDF; modo "1 cartao por pagina" recorta cada cartao.
"""

import io
import re
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

from src.utils.paths import get_templates_dir, get_certificados_dir
from src.utils.validators import formatar_telefone

TOKEN_RE = re.compile(r"\{\{([A-Z_]+)\}\}")
SLOT_NAME_RE = re.compile(r"^CARD(\d+)_")
PHOTO_SUFFIX = "_FOTO"

EMPLOYEE_FIELDS = {"NOME", "FUNCAO", "TELEFONE", "CPF", "MATRICULA"}
BATCH_FIELDS = {"SETOR", "PAPEL"}
TEMPLATE_FIELDS = {"EMPRESA"}
ALL_FIELDS = EMPLOYEE_FIELDS | BATCH_FIELDS | TEMPLATE_FIELDS

PP_SAVE_AS_PDF = 32


# ── Helpers de navegacao de shapes ──────────────────────────

def _iter_all_shapes(shapes):
    """Itera shapes recursivamente (entra em grupos)."""
    for shp in shapes:
        yield shp
        if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
            try:
                yield from _iter_all_shapes(shp.shapes)
            except Exception:
                pass


def _shape_slot(shape) -> Optional[int]:
    """Retorna o numero do slot (1-based) se o shape for nomeado CARD{k}_..."""
    m = SLOT_NAME_RE.match(shape.name or "")
    return int(m.group(1)) if m else None


# ── Carga e deteccao de campos ──────────────────────────────

def load_pptx_card_templates() -> Dict[str, dict]:
    """Carrega templates/cards/pptx/*.card.json e injeta used_fields/uses_photo/capacity."""
    pptx_dir = get_templates_dir() / "cards" / "pptx"
    templates: Dict[str, dict] = {}
    if not pptx_dir.exists():
        return templates
    for f in sorted(pptx_dir.glob("*.card.json")):
        try:
            import json
            data = json.loads(f.read_text(encoding="utf-8"))
            if data.get("card_code") and data.get("template_type") == "pptx":
                pptx_path = pptx_dir / data.get("pptx_file", "")
                if pptx_path.exists():
                    data["_pptx_path"] = str(pptx_path)
                    used, uses_photo = detect_used_fields(pptx_path)
                    data["used_fields"] = sorted(used)
                    data["uses_photo"] = uses_photo
                    try:
                        prs = Presentation(str(pptx_path))
                        data["_slides"] = len(prs.slides)
                        data["_capacity"] = len(prs.slides) * int(data.get("cards_per_slide", 1))
                    except Exception:
                        data["_slides"] = int(data.get("slides", 1))
                        data["_capacity"] = data["_slides"] * int(data.get("cards_per_slide", 1))
                    templates[data["card_code"]] = data
        except Exception:
            continue
    return templates


def detect_used_fields(pptx_path: Path) -> Tuple[set, bool]:
    """Escaneia o PPTX: tokens {{CAMPO}} usados e se ha shapes de foto (CARD*_FOTO)."""
    used = set()
    uses_photo = False
    try:
        prs = Presentation(str(pptx_path))
    except Exception:
        return used, uses_photo
    try:
        for slide in prs.slides:
            for shape in _iter_all_shapes(slide.shapes):
                name = shape.name or ""
                if name.endswith(PHOTO_SUFFIX) and SLOT_NAME_RE.match(name):
                    uses_photo = True
                if getattr(shape, "has_text_frame", False):
                    try:
                        text = shape.text_frame.text
                    except Exception:
                        continue
                    for m in TOKEN_RE.finditer(text):
                        field = m.group(1)
                        if field in ALL_FIELDS:
                            used.add(field)
    finally:
        pass
    return used, uses_photo


def pptx_capacity(template: dict) -> int:
    return int(template.get("_capacity") or 1)


# ── Validacao dinamica ──────────────────────────────────────

def validate_employees_for_pptx(employees: list, template: dict) -> Tuple[list, List[str]]:
    """
    Validacao conforme os campos que o template realmente usa:
    - foto: exigida se o template tem shapes CARD*_FOTO
    - telefone: exigido se {{TELEFONE}} esta no template
    - matricula: nunca bloqueia (fallback CPF)
    """
    used = set(template.get("used_fields") or [])
    needs_photo = bool(template.get("uses_photo"))
    needs_tel = "TELEFONE" in used
    valid, missing = [], []
    for emp in employees:
        faltas = []
        if needs_tel and not getattr(emp, "telefone", None):
            faltas.append("telefone")
        if needs_photo and not getattr(emp, "foto", None):
            faltas.append("foto 3x4")
        if faltas:
            missing.append(f"{emp.nome}: sem {' e '.join(faltas)}")
        else:
            valid.append(emp)
    return valid, missing


# ── Substituicao de texto ───────────────────────────────────

_FONT_CACHE: Dict[str, Optional[str]] = {}


def _find_font_path(bold: bool) -> Optional[str]:
    """Localiza fonte para medir largura de texto (Segoe UI -> DejaVu -> None)."""
    key = "bold" if bold else "normal"
    if key in _FONT_CACHE:
        return _FONT_CACHE[key]
    import os

    path = None
    win_dir = os.environ.get("WINDIR", r"C:\Windows")
    for name in (["segoeuib.ttf", "seguisb.ttf"] if bold else ["segoeui.ttf", "segoeuisl.ttf"]):
        cand = Path(win_dir) / "Fonts" / name
        if cand.exists():
            path = str(cand)
            break
    if path is None:
        dv = get_templates_dir().parent / "assets" / "fonts" / (
            "DejaVuSans-Bold.ttf" if bold else "DejaVuSans.ttf")
        if dv.exists():
            path = str(dv)
    _FONT_CACHE[key] = path
    return path


def _text_width_pt(text: str, size_pt: float, bold: bool) -> float:
    """Largura aproximada do texto em pontos (medicao PIL ou heuristica)."""
    path = _find_font_path(bold)
    if path:
        try:
            from PIL import ImageFont

            f = ImageFont.truetype(path, int(round(size_pt * 4)))
            return f.getlength(text) / 4.0
        except Exception:
            pass
    return len(text) * size_pt * (0.62 if bold else 0.58)


LINE_HEIGHT_FACTOR = 1.28  # altura de linha estimada (x fonte) para o shrink vertical


def _wrap_line_count(text: str, size_pt: float, bold: bool, usable_pt: float) -> int:
    """Estima o numero de linhas com quebra gulosa de palavras (medindo cada linha)."""
    words = [w for w in text.split() if w]
    if not words:
        return 1
    lines, cur = 1, ""
    for w in words:
        cand = f"{cur} {w}".strip()
        if not cur or _text_width_pt(cand, size_pt, bold) <= usable_pt:
            cur = cand
        else:
            lines += 1
            cur = w
    return lines


def _usable_width_pt(shape, text_frame) -> float:
    usable_emu = (shape.width or 0) - (text_frame.margin_left or 0) - (text_frame.margin_right or 0)
    return max(usable_emu, 0) / 12700.0


def _shrink_paragraph_to_fit(para, text: str, shape, text_frame):
    """
    Reduz a fonte do 1o run ate o texto caber no shape, considerando a quebra
    de linha (word_wrap): largura da maior palavra E altura do bloco de linhas.
    Passos de 0.5pt, minimo 5.5pt.
    """
    runs = para.runs
    if not runs:
        return
    r0 = runs[0]
    size = r0.font.size.pt if r0.font.size else 12.0
    if size <= 5.5:
        return
    uw_pt = _usable_width_pt(shape, text_frame)
    if uw_pt <= 0:
        return
    uh_emu = (shape.height or 0) - (text_frame.margin_top or 0) - (text_frame.margin_bottom or 0)
    uh_pt = uh_emu / 12700.0 if uh_emu > 0 else None
    bold = bool(r0.font.bold)
    words = [w for w in text.split() if w]

    new_size = size
    while new_size > 5.5:
        longest_word = max((_text_width_pt(w, new_size, bold) for w in words), default=0.0)
        n_lines = _wrap_line_count(text, new_size, bold, uw_pt)
        fits_w = longest_word <= uw_pt
        fits_h = uh_pt is None or (n_lines * new_size * LINE_HEIGHT_FACTOR) <= uh_pt
        if fits_w and fits_h:
            break
        new_size -= 0.5
    new_size = max(5.5, new_size)
    if new_size < size:
        from pptx.util import Pt

        r0.font.size = Pt(new_size)


def _clip_text_to_width(text: str, size_pt: float, bold: bool, usable_pt: float) -> str:
    """Modo clip (sem quebra): corta caracteres ate caber na largura util."""
    if _text_width_pt(text, size_pt, bold) <= usable_pt:
        return text
    out = text
    while out and _text_width_pt(out, size_pt, bold) > usable_pt:
        out = out[:-1].rstrip()
    return out


def _replace_tokens_in_text_frame(shape, values: Dict[str, str], fit_mode: str = "wrap"):
    """
    Substitui {{TOKENS}} a nivel de paragrafo (resolve runs quebrados tipo
    "FABRICIO " + "CARVALHO"). O texto substituido herda o estilo do 1o run.

    fit_mode:
    - "wrap" (default): liga word_wrap e reduz a fonte ate as linhas quebradas
      caberem na largura (maior palavra) e altura do shape
    - "clip": sem quebra e sem reducao — corta o que passar do limite do campo
    """
    text_frame = shape.text_frame
    try:
        text_frame.word_wrap = (fit_mode != "clip")
        if fit_mode != "clip":
            # caixas desses templates tem altura exata de 1 linha; zerar as
            # margens verticais libera espaco para a 2a linha da quebra
            text_frame.margin_top = 0
            text_frame.margin_bottom = 0
    except Exception:
        pass
    for para in text_frame.paragraphs:
        runs = para.runs
        if not runs:
            continue
        full = "".join(r.text for r in runs)
        if not TOKEN_RE.search(full):
            continue
        new_text = TOKEN_RE.sub(lambda m: values.get(m.group(1), m.group(0)), full)
        if fit_mode == "clip":
            r0 = runs[0]
            size = r0.font.size.pt if r0.font.size else 12.0
            new_text = _clip_text_to_width(new_text, size, bool(r0.font.bold),
                                           _usable_width_pt(shape, text_frame))
        runs[0].text = new_text
        for r in runs[1:]:
            r.text = ""
        if fit_mode != "clip":
            try:
                _shrink_paragraph_to_fit(para, new_text, shape, text_frame)
            except Exception:
                pass


# ── Substituicao de foto ────────────────────────────────────

def _crop_photo_to_box(foto_bytes: bytes, box_w_emu: int, box_h_emu: int) -> bytes:
    """Center-crop da foto para o aspecto do box (evita distorcao). Retorna JPEG."""
    from PIL import Image

    img = Image.open(io.BytesIO(foto_bytes))
    if img.mode in ("RGBA", "LA", "P"):
        img = img.convert("RGB")
    target_ratio = box_w_emu / max(box_h_emu, 1)
    w, h = img.size
    cur_ratio = w / max(h, 1)
    if abs(cur_ratio - target_ratio) > 0.01:
        if cur_ratio > target_ratio:
            new_w = int(h * target_ratio)
            x0 = (w - new_w) // 2
            img = img.crop((x0, 0, x0 + new_w, h))
        else:
            new_h = int(w / target_ratio)
            y0 = (h - new_h) // 2
            img = img.crop((0, y0, w, y0 + new_h))
    out = io.BytesIO()
    img.save(out, format="JPEG", quality=88)
    return out.getvalue()


def _gray_placeholder_jpeg(w: int = 120, h: int = 160) -> bytes:
    from PIL import Image

    img = Image.new("RGB", (w, h), (224, 224, 224))
    out = io.BytesIO()
    img.save(out, format="JPEG")
    return out.getvalue()


def _swap_picture_image(pic_shape, img_bytes: bytes):
    """Troca a imagem de um shape Picture mantendo posicao/tamanho/z-order."""
    image_part, rId = pic_shape.part.get_or_add_image_part(io.BytesIO(img_bytes))
    blip = pic_shape._element.blipFill.blip
    blip.set(qn("r:embed"), rId)


def _try_swap_photo(shape, foto_bytes: Optional[bytes]):
    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        return False
    try:
        if foto_bytes:
            data = _crop_photo_to_box(foto_bytes, shape.width or 1, shape.height or 1)
        else:
            data = _gray_placeholder_jpeg()
        _swap_picture_image(shape, data)
        return True
    except Exception:
        return False


# ── Preenchimento de um clone ───────────────────────────────

def _employee_values(emp, template: dict, options: dict) -> Dict[str, str]:
    papeis = options.get("papeis") or {}
    # matricula: exclusiva da emissao (numero tem validade) — vem so do popup, sem fallback
    matricula = (options.get("matriculas", {}).get(emp.id) or "").strip()
    tel = formatar_telefone(emp.telefone) if getattr(emp, "telefone", None) else "-"
    return {
        "NOME": (emp.nome or "").upper(),
        "FUNCAO": (emp.funcao or "-").upper(),
        "TELEFONE": tel,
        "CPF": emp.cpf or "",
        "MATRICULA": matricula,
        "SETOR": options.get("setor", ""),
        "EMPRESA": template.get("empresa_default", "ALTEC"),
        "PAPEL": papeis.get(emp.id, "LIDERADO"),
    }


def _blank_values(template: dict) -> Dict[str, str]:
    return {f: "" for f in (template.get("used_fields") or [])}


def _fill_clone(prs, chunk: list, template: dict, options: dict):
    """
    Preenche um clone do template com ate `capacity` funcionarios.
    Slots nao usados ficam em branco (texto vazio + foto cinza).
    """
    cards_per_slide = int(template.get("cards_per_slide", 1))
    fit_mode = template.get("text_fit", "wrap")
    slides = list(prs.slides)

    for idx in range(len(slides) * cards_per_slide):
        emp = chunk[idx] if idx < len(chunk) else None
        slide_idx = idx // cards_per_slide
        slot = idx % cards_per_slide + 1
        if slide_idx >= len(slides):
            break
        values = _employee_values(emp, template, options) if emp else _blank_values(template)
        foto = emp.foto if (emp and getattr(emp, "foto", None)) else None

        for shape in _iter_all_shapes(slides[slide_idx].shapes):
            if _shape_slot(shape) != slot:
                continue
            if getattr(shape, "has_text_frame", False):
                try:
                    _replace_tokens_in_text_frame(shape, values, fit_mode)
                except Exception:
                    pass
            if (shape.name or "").endswith(PHOTO_SUFFIX):
                _try_swap_photo(shape, foto)


# ── Conversao PPTX -> PDF (PowerPoint via COM) ──────────────

def _pptx_to_pdf_batch(pairs: List[Tuple[Path, Path]]):
    """
    Converte varios PPTX em PDF usando UMA sessao do PowerPoint.
    pairs: [(pptx_path, pdf_path), ...]
    """
    try:
        import comtypes.client
    except ImportError:
        raise RuntimeError("Biblioteca comtypes nao instalada. Instale com: pip install comtypes")

    try:
        app = comtypes.client.CreateObject("PowerPoint.Application", dynamic=True)
    except Exception as e:
        raise RuntimeError(
            "Microsoft PowerPoint nao encontrado neste computador.\n"
            "Templates PPTX exigem o PowerPoint instalado (Office).\n"
            f"Detalhe: {e}"
        ) from e

    try:
        for src, dst in pairs:
            pres = app.Presentations.Open(str(Path(src).resolve()), True, False, False)
            try:
                pres.SaveAs(str(Path(dst).resolve()), PP_SAVE_AS_PDF)
            finally:
                pres.Close()
    finally:
        try:
            app.Quit()
        except Exception:
            pass


# ── Merge e recorte de PDFs (PyMuPDF) ───────────────────────

def _merge_pdfs(src_paths: List[Path], out_path: Path):
    import fitz

    doc = fitz.open()
    try:
        for p in src_paths:
            src = fitz.open(str(p))
            try:
                doc.insert_pdf(src)
            finally:
                src.close()
        doc.save(str(out_path))
    finally:
        doc.close()


def _split_one_per_page(
    sheet_pdfs: List[Path],
    employees: list,
    template: dict,
    out_path: Path,
    only_employee: Optional[object] = None,
) -> Optional[Path]:
    """
    Recorta cada cartao do(s) PDF(s) de folha em uma pagina propria.
    Mapeia: funcionario global g -> folha j=(g//C), posicao idx=(g%C),
    slide s=idx//K, slot k=idx%K -> pagina (j*S + s), zona zones[k].
    Se only_employee, gera apenas o cartao dele em out_path.
    """
    import fitz

    cards_per_slide = int(template.get("cards_per_slide", 1))
    zones = template.get("card_zones_fraction") or []
    if len(zones) < cards_per_slide:
        zones = zones + [[0.0, 0.0, 1.0, 1.0]] * (cards_per_slide - len(zones))
    capacity = int(template.get("_capacity") or cards_per_slide)
    slides_per_clone = max(1, capacity // cards_per_slide)

    src_doc = fitz.open()
    try:
        for p in sheet_pdfs:
            d = fitz.open(str(p))
            src_doc.insert_pdf(d)
            d.close()

        out = fitz.open()
        try:
            indices = range(len(employees))
            if only_employee is not None:
                indices = [i for i, e in enumerate(employees) if e is only_employee]
                if not indices:
                    return None
            for g in indices:
                j, idx = divmod(g, capacity)
                s, k = divmod(idx, cards_per_slide)
                page_no = j * slides_per_clone + s
                if page_no >= len(src_doc):
                    continue
                page = src_doc[page_no]
                pw, ph = page.rect.width, page.rect.height
                z = zones[k]
                clip = fitz.Rect(z[0] * pw, z[1] * ph, z[2] * pw, z[3] * ph)
                newpage = out.new_page(width=clip.width, height=clip.height)
                newpage.show_pdf_page(newpage.rect, src_doc, page_no, clip=clip)
            out.save(str(out_path))
        finally:
            out.close()
        return out_path
    finally:
        src_doc.close()


# ── Geracao principal ───────────────────────────────────────

def _sanitize_filename(name: str) -> str:
    s = re.sub(r'[^\w\s-]', '', name).strip()
    return re.sub(r'[\s]+', '_', s)[:50]


def generate_pptx_cards(
    employees: list,
    template: dict,
    options: Optional[dict] = None,
    single_pdf: bool = True,
    one_per_page: bool = False,
    output_dir: Optional[Path] = None,
) -> Tuple[List[Path], List[str]]:
    """
    Gera cartoes de bloqueio a partir de template PPTX.

    - single_pdf=True: um PDF unico (folhas do template ou 1 cartao/pagina)
    - single_pdf=False: um PDF por funcionario (cartao recortado)

    options: {"setor": str, "papeis": {employee_id: "LIDER"|"LIDERADO"}}
    Retorna (caminhos_gerados, faltantes_msg).
    """
    options = options or {}
    valid, missing = validate_employees_for_pptx(employees, template)
    if not valid:
        return [], missing

    card_code = template.get("card_code", "PPTX")
    cliente_dir = output_dir or (get_certificados_dir() / "CARTOES" / card_code)
    cliente_dir.mkdir(parents=True, exist_ok=True)

    pptx_path = Path(template["_pptx_path"])
    capacity = pptx_capacity(template)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")

    # 1. Criar clones preenchidos (chunks de `capacity`)
    import tempfile

    tmp_dir = Path(tempfile.mkdtemp(prefix="cartoes_pptx_"))
    clone_pptx: List[Path] = []
    try:
        chunks = [valid[i:i + capacity] for i in range(0, len(valid), capacity)]
        for ci, chunk in enumerate(chunks):
            prs = Presentation(str(pptx_path))
            _fill_clone(prs, chunk, template, options)
            clone_path = tmp_dir / f"lote_{ci + 1}.pptx"
            prs.save(str(clone_path))
            clone_pptx.append(clone_path)

        # 2. Converter todos os clones em PDF (uma sessao do PowerPoint)
        sheet_pdfs = [clone_path.with_suffix(".pdf") for clone_path in clone_pptx]
        _pptx_to_pdf_batch(list(zip(clone_pptx, sheet_pdfs)))

        # 3. Montar saida conforme modo
        generated: List[Path] = []
        if single_pdf:
            out = cliente_dir / f"CARTOES_{card_code}_{timestamp}.pdf"
            if one_per_page:
                _split_one_per_page(sheet_pdfs, valid, template, out)
            else:
                _merge_pdfs(sheet_pdfs, out)
            generated.append(out)
        else:
            for emp in valid:
                out = cliente_dir / f"CARTAO_{_sanitize_filename(emp.nome)}_{card_code}.pdf"
                _split_one_per_page(sheet_pdfs, valid, template, out, only_employee=emp)
                generated.append(out)

        return generated, missing
    finally:
        import shutil

        shutil.rmtree(tmp_dir, ignore_errors=True)
