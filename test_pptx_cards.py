"""
Testes do servico de cartoes PPTX.

1. Testes unitarios (sem PowerPoint): tokens, validacao, fallback matricula
2. Teste E2E (requer PowerPoint): gera PDFs reais para cada template

Uso:  python test_pptx_cards.py [--unit] [--e2e]
"""
import io
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parent
sys.path.insert(0, str(ROOT))

from src.core.models import Employee
from src.core.pptx_card_service import (
    load_pptx_card_templates,
    validate_employees_for_pptx,
    generate_pptx_cards,
    _replace_tokens_in_text_frame,
    _employee_values,
)
from pptx import Presentation


def make_employees():
    fotos_dir = ROOT / "FOTOS PARA TESTE"
    foto1 = next(fotos_dir.glob("ADELSON*")).read_bytes()
    foto2 = next(fotos_dir.glob("ALCIMAR*")).read_bytes()
    fotos = [foto1, foto2]
    base = [
        # nome longo de proposito: estressa o auto-shrink de fonte
        ("Fabricio Carvalho Ferreira Da Silva Santos Junior", "529.982.247-25", "Soldador", "21981586650"),
        ("Ailton Costa Martins", "111.444.777-35", "Tecnico em Seguranca do Trabalho", "22981632680"),
        ("Luciano Dias", "123.456.789-09", "Encarregado de Montagem e Manutencao Predial", "21984209236"),
        ("Paulo Vitor Barbosa", "987.654.321-00", "Eletricista", "21987654321"),
        ("Herkton Diniz Santos", "222.333.444-05", "Mecanico de Manutencao Industrial", "22987651234"),
        ("Valques Ribeiro Goncalves", "333.444.555-08", "Caldeireiro", "21986543210"),
        ("Wesley De Azevedo Barboza", "444.555.666-19", "Montador de Andaime", "22985432109"),
        ("Wancler Alves Poubel", "555.666.777-20", "Ajudante de Soldador Pleno", "21984321098"),
    ]
    return [
        Employee(id=i + 1, nome=n, cpf=c, funcao=f, telefone=t, foto=fotos[i % 2])
        for i, (n, c, f, t) in enumerate(base)
    ]


E2E_MATRICULAS = {
    1: "12240721707", 2: "22927112738", 3: "10442439709", 4: "70132975",
    5: "13762726795", 6: "14154086707", 7: "17088417766", 8: "18393107741",
}


# ── Unitarios ───────────────────────────────────────────────

def test_tokens():
    class FakeRun:
        def __init__(self, t):
            self.text = t

    class FakePara:
        def __init__(self, *ts):
            self.runs = [FakeRun(t) for t in ts]

    class FakeTF:
        def __init__(self, paras):
            self.paragraphs = paras
            self.margin_left = 91440
            self.margin_right = 91440
            self.margin_top = 45720
            self.margin_bottom = 45720

    class FakeShape:
        def __init__(self, tf):
            self.text_frame = tf
            self.width = 3600000  # 10cm — largo o bastante para nao acionar shrink
            self.height = None

    # placeholder quebrado em varios runs
    tf = FakeTF([FakePara("Nome: ", "{{NO", "ME}} | teste")])
    _replace_tokens_in_text_frame(FakeShape(tf), {"NOME": "FABIO"})
    assert tf.paragraphs[0].runs[0].text == "Nome: FABIO | teste", tf.paragraphs[0].runs[0].text
    assert all(r.text == "" for r in tf.paragraphs[0].runs[1:])
    print("[OK] token split em runs")


def test_shrink():
    class FakeRun:
        def __init__(self, t):
            self.text = t
            self.font = type("F", (), {"size": type("S", (), {"pt": 12.0})(), "bold": True})()

    class FakePara:
        def __init__(self, r):
            self.runs = [r]

    class FakeTF:
        def __init__(self, para):
            self.paragraphs = [para]
            self.margin_left = 91440
            self.margin_right = 91440
            self.margin_top = 45720
            self.margin_bottom = 45720

    class FakeShape:
        def __init__(self, tf, w, h=None):
            self.text_frame = tf
            self.width = w
            self.height = h  # None -> shrink ignora altura (so largura)

    from src.core.pptx_card_service import _shrink_paragraph_to_fit

    # texto longo em shape estreito (2cm), sem altura -> encolhe pela largura
    r = FakeRun("FABRICIO CARVALHO FERREIRA DA SILVA SANTOS JUNIOR")
    _shrink_paragraph_to_fit(FakePara(r), r.text, FakeShape(None, 720000), FakeTF(None))
    assert r.font.size.pt < 12.0, r.font.size.pt
    print(f"[OK] auto-shrink (largura): 12pt -> {r.font.size.pt}pt para caber em 2cm")

    # com quebra: caixa larga e alta comporta em 2 linhas sem reduzir tanto
    r2 = FakeRun("FABRICIO CARVALHO FERREIRA DA SILVA")
    _shrink_paragraph_to_fit(FakePara(r2), r2.text, FakeShape(None, 2000000, 1000000), FakeTF(None))
    print(f"[OK] auto-shrink (wrap vertical): {r2.font.size.pt}pt em caixa 5.6x2.8cm")


def test_matricula_do_popup():
    emps = make_employees()
    tpl = {"empresa_default": "ALTEC"}
    # matricula vem exclusivamente do popup (sem fallback CPF)
    vals = _employee_values(emps[0], tpl, {"matriculas": {emps[0].id: "99988877766"}})
    assert vals["MATRICULA"] == "99988877766", vals["MATRICULA"]
    # sem informar no popup -> vazio (UI bloqueia esse caso antes de chegar aqui)
    vals2 = _employee_values(emps[1], tpl, {})
    assert vals2["MATRICULA"] == "", vals2["MATRICULA"]
    print("[OK] matricula exclusiva do popup (sem fallback CPF)")


def test_wrap_e_clip():
    from src.core.pptx_card_service import _wrap_line_count, _clip_text_to_width, _text_width_pt

    # wrap: nome longo em caixa estreita -> 2+ linhas
    nome = "FABRICIO CARVALHO FERREIRA DA SILVA SANTOS JUNIOR"
    largura_1linha = _text_width_pt(nome, 8.0, True)
    n = _wrap_line_count(nome, 8.0, True, largura_1linha * 0.45)
    assert n >= 2, n
    n1 = _wrap_line_count("JOAO", 8.0, True, largura_1linha)
    assert n1 == 1
    print(f"[OK] estimador de quebra: {n} linhas para caixa de 45%")

    # clip: texto cortado no limite, prefixo preservado, sem reticencias
    limit = _text_width_pt("FABRICIO CARVALHO", 8.0, True)
    cortado = _clip_text_to_width("FABRICIO CARVALHO FERREIRA", 8.0, True, limit)
    assert cortado.startswith("FABRICIO"), cortado
    assert "..." not in cortado, cortado
    assert _text_width_pt(cortado, 8.0, True) <= limit, cortado
    print(f"[OK] clip: 'FABRICIO CARVALHO FERREIRA' -> '{cortado}'")


def test_clip_linhas_csn():
    """Modo CSN: excedente da 2a linha do nome e da 1a linha dos demais e apagado."""
    from src.core.pptx_card_service import _clip_text_to_lines, _text_width_pt

    texto = "UM DOIS TRES QUATRO CINCO SEIS SETE"
    # largura que comporta o par mais largo ("TRES QUATRO") por linha
    w2 = _text_width_pt("TRES QUATRO", 10.0, False)
    duas = _clip_text_to_lines(texto, 10.0, False, w2, 2)
    uma = _clip_text_to_lines(texto, 10.0, False, w2, 1)
    # caixa comporta 2 palavras por linha: 2 linhas = 4 palavras; 1 linha = 2
    assert duas.split() == ["UM", "DOIS", "TRES", "QUATRO"], duas
    assert uma.split() == ["UM", "DOIS"], uma
    print(f"[OK] clip por linhas: 2 linhas -> '{duas}' | 1 linha -> '{uma}'")


def test_edicao_por_emissao():
    """Copias transitorias: edicao nao afeta o funcionario original."""
    orig = make_employees()[0]
    copia = orig.model_copy()
    copia.nome = "EDITADO SO NA EMISSAO"
    copia.funcao = "Nova Funcao"
    vals = _employee_values(copia, {"empresa_default": "ALTEC"}, {})
    assert vals["NOME"] == "EDITADO SO NA EMISSAO"
    assert vals["FUNCAO"] == "NOVA FUNCAO"
    assert orig.nome != "EDITADO SO NA EMISSAO" and orig.funcao == "Soldador"
    assert copia.id == orig.id  # papeis/matriculas continuam casando por id
    print("[OK] edicao por emissao: copia editada, original intacto")


def test_validacao_dinamica():
    emps = make_employees()
    tpls = load_pptx_card_templates()

    # sem telefone e sem foto
    sem_tudo = Employee(id=9, nome="Sem Dados", cpf=None)
    valid, missing = validate_employees_for_pptx([sem_tudo], tpls["ARCELORMITTAL"])
    # ARCELORMITTAL nao usa telefone -> so cobra foto
    assert not valid and "foto" in missing[0] and "telefone" not in missing[0], missing
    print("[OK] ARCELORMITTAL exige foto mas nao telefone")

    valid, missing = validate_employees_for_pptx([sem_tudo], tpls["ALTEC-PEQUENO"])
    # ALTEC-PEQUENO usa telefone e foto (placeholders adicionados na preparacao)
    assert not valid and "telefone" in missing[0] and "foto" in missing[0], missing
    print("[OK] ALTEC-PEQUENO exige telefone e foto")

    valid, missing = validate_employees_for_pptx(emps, tpls["CSN"])
    assert len(valid) == len(emps) and not missing
    print("[OK] funcionarios completos validos")


def test_fill_clone_tokens():
    tpls = load_pptx_card_templates()
    tpl = tpls["LOTOTO"]
    assert tpl.get("text_fit") == "clip", "LOTOTO deveria usar text_fit=clip"
    prs = Presentation(tpl["_pptx_path"])
    from src.core.pptx_card_service import _fill_clone
    emps = make_employees()[:2]
    _fill_clone(prs, emps, tpl, {"setor": "Manutencao Mecanica", "papeis": {emps[0].id: "LIDER"},
                                  "matriculas": E2E_MATRICULAS})
    texts = []
    from src.core.pptx_card_service import _iter_all_shapes
    for slide in prs.slides:
        for shp in _iter_all_shapes(slide.shapes):
            if getattr(shp, "has_text_frame", False):
                texts.append(shp.text_frame.text)
                # modo clip: word_wrap desligado nos shapes de dados preenchidos
                if (shp.name or "").startswith("CARD"):
                    assert shp.text_frame.word_wrap is False, shp.name
    joined = "\n".join(texts)
    assert "{{" not in joined, "tokens remanescentes!"
    assert "FABRICIO CARVALHO" in joined  # clip pode cortar o fim, nunca o inicio
    assert "LIDER" in joined and "Manutencao Mecanica" in joined
    print("[OK] _fill_clone LOTOTO (clip, 2 slides) sem tokens remanescentes")


def test_altec_pequeno_8_slots():
    """ALTEC-PEQUENO: 8 slots com o 8o independente do 7o (fix do slot 8)."""
    tpls = load_pptx_card_templates()
    tpl = tpls["ALTEC-PEQUENO"]
    assert tpl["_capacity"] == 8, tpl["_capacity"]
    prs = Presentation(tpl["_pptx_path"])
    from src.core.pptx_card_service import _fill_clone, _iter_all_shapes
    emps = make_employees()
    _fill_clone(prs, emps, tpl, {})
    texts = [s.text_frame.text for s in _iter_all_shapes(prs.slides[0].shapes)
             if getattr(s, "has_text_frame", False)]
    joined = "\n".join(texts)
    assert "WESLEY DE AZEVEDO BARBOZA" in joined, "slot 7 sem dados"
    assert "WANCLER ALVES POUBEL" in joined, "slot 8 sem dados (renomeacao falhou?)"
    print("[OK] ALTEC-PEQUENO: 8 slots preenchidos, 8o distinto do 7o")


# ── E2E ─────────────────────────────────────────────────────

def _check_data_inside_zones(pdf_path, tpl, emps, options):
    """Verifica que todo texto de DADOS preenchido esta dentro da zona do cartao."""
    import fitz

    needles = set()
    for e in emps:
        vals = _employee_values(e, tpl, options)
        needles.update(v for v in vals.values() if isinstance(v, str) and len(v) >= 4)
    needles |= {"LIDER", "LIDERADO"}

    bad = []
    doc = fitz.open(str(pdf_path))
    try:
        for pno in range(len(doc)):
            page = doc[pno]
            pw, ph = page.rect.width, page.rect.height
            zones = [fitz.Rect(z[0] * pw, z[1] * ph, z[2] * pw, z[3] * ph)
                     for z in tpl["card_zones_fraction"]]
            for b in page.get_text("dict")["blocks"]:
                if b["type"] != 0:
                    continue
                for l in b["lines"]:
                    for s in l["spans"]:
                        t = s["text"].strip()
                        if not t or not any(n in t for n in needles):
                            continue
                        r = fitz.Rect(s["bbox"])
                        c = fitz.Point((r.x0 + r.x1) / 2, (r.y0 + r.y1) / 2)
                        z = next((z for z in zones if z.contains(c)), None)
                        if z is None:
                            bad.append((pno, t[:45], "fora de qualquer cartao"))
                            continue
                        tol = 3.0
                        zt = fitz.Rect(z.x0 - tol, z.y0 - tol, z.x1 + tol, z.y1 + tol)
                        if not zt.contains(r):
                            bad.append((pno, t[:45], "estoura a borda do cartao"))
    finally:
        doc.close()
    return bad


def test_e2e():
    tpls = load_pptx_card_templates()
    emps = make_employees()
    options = {"setor": "Manutencao Mecanica", "papeis": {1: "LIDER"},
               "matriculas": E2E_MATRICULAS}
    out = ROOT / "CERTIFICADOS" / "CARTOES" / "_TESTE"
    for code, tpl in tpls.items():
        paths, missing = generate_pptx_cards(emps, tpl, options=options,
                                             single_pdf=True, one_per_page=False,
                                             output_dir=out)
        assert paths, f"{code}: nenhum PDF gerado ({missing})"
        p = paths[0]
        assert p.exists() and p.stat().st_size > 10_000, f"{code}: PDF muito pequeno"
        print(f"[OK] {code}: {p.name} ({p.stat().st_size // 1024} KB)")

        # dados preenchidos precisam caber dentro do cartao (auto-shrink)
        bad = _check_data_inside_zones(p, tpl, emps, options)
        assert not bad, f"{code}: dados fora do cartao: {bad[:5]}"
        print(f"[OK] {code}: dados dentro dos limites dos cartoes")

        # modo 1 cartao por pagina
        paths1, _ = generate_pptx_cards(emps, tpl, options=options,
                                        single_pdf=True, one_per_page=True,
                                        output_dir=out)
        p1 = paths1[0]
        import fitz
        d = fitz.open(str(p1))
        try:
            expected = len(emps)  # 1 pagina por cartao
            assert len(d) == expected, f"{code}: esperava {expected} paginas, veio {len(d)}"
        finally:
            d.close()
        print(f"[OK] {code} 1/pagina: {p1.name} ({len(fitz.open(str(p1)))} paginas)")


if __name__ == "__main__":
    args = set(sys.argv[1:])
    run_all = not args
    if run_all or "--unit" in args:
        test_tokens()
        test_shrink()
        test_wrap_e_clip()
        test_clip_linhas_csn()
        test_edicao_por_emissao()
        test_matricula_do_popup()
        test_validacao_dinamica()
        test_fill_clone_tokens()
        test_altec_pequeno_8_slots()
    if run_all or "--e2e" in args:
        test_e2e()
    print("\nTODOS OS TESTES PASSARAM")
